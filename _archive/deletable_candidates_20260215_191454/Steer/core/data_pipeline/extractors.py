# extractors.py - Extracted from pipeline.py (Document text extractors)
"""Document text extraction for various file formats."""

from __future__ import annotations

import io
import os
import platform
from pathlib import Path
from typing import Any, Dict, List, Tuple

from core.utils.nlp import TextCleaner

# Optional dependencies
try:
    import pandas as pd
except Exception:
    pd = None

try:
    import docx
except Exception:
    docx = None

try:
    import pptx
except Exception:
    pptx = None

try:
    import fitz  # PyMuPDF
except Exception:
    fitz = None

try:
    import pdfplumber
except Exception:
    pdfplumber = None

try:
    from pdfminer.high_level import extract_text as pdfminer_extract_text
except Exception:
    pdfminer_extract_text = None

try:
    import olefile
except Exception:
    olefile = None

try:
    import pyhwp
except Exception:
    pyhwp = None

try:
    import win32com.client as win32com
    import pythoncom
except Exception:
    win32com = None
    pythoncom = None


class BaseExtractor:
    exts: Tuple[str, ...] = ()

    def can_handle(self, path: Path) -> bool:
        return path.suffix.lower() in self.exts

    def extract(self, path: Path) -> Dict[str, Any]:
        raise NotImplementedError

class HwpExtractor(BaseExtractor):
    exts = (".hwp",)

    def extract(self, p: Path) -> Dict[str, Any]:
        system = platform.system().lower()
        if system.startswith("win") and win32com:
            com_initialized = False
            try:
                if pythoncom:
                    pythoncom.CoInitialize()
                    com_initialized = True
                app = win32com.Dispatch("HWPFrame.HwpObject")
                try:
                    app.Open(str(p))
                    text = app.GetTextFile("TEXT", "") or ""
                    return {
                        "ok": True,
                        "text": TextCleaner.clean(text),
                        "meta": {"engine": "win32com-hwp"},
                    }
                finally:
                    try:
                        app.Quit()
                    except Exception:
                        pass
            except Exception as exc:
                return {"ok": False, "text": "", "meta": {"error": f"HWP win32com 실패: {exc}"}}
            finally:
                if com_initialized and pythoncom:
                    try:
                        pythoncom.CoUninitialize()
                    except Exception:
                        pass
        if olefile and pyhwp:
            try:
                from pyhwp.hwp5txt import hwp5txt  # type: ignore

                with olefile.OleFileIO(str(p)) as ole:
                    buf = io.StringIO()
                    hwp5txt(ole, buf)
                    text = buf.getvalue()
                cleaned = TextCleaner.clean(text)
                if cleaned:
                    return {
                        "ok": True,
                        "text": cleaned,
                        "meta": {"engine": "pyhwp", "bytes": p.stat().st_size},
                    }
            except Exception as exc:
                return {
                    "ok": False,
                    "text": "",
                    "meta": {"error": f"HWP pyhwp 추출 실패: {exc}"},
                }
        return {
            "ok": False,
            "text": "",
            "meta": {"error": "HWP 추출을 위해서는 Windows + 한/글 환경이 필요합니다."},
        }


class DocDocxExtractor(BaseExtractor):
    exts = (".doc", ".docx")

    def extract(self, p: Path) -> Dict[str, Any]:
        suffix = p.suffix.lower()
        if suffix == ".docx" and docx:
            try:
                document = docx.Document(str(p))
                text = "\n".join(par.text for par in document.paragraphs)
                return {
                    "ok": True,
                    "text": TextCleaner.clean(text),
                    "meta": {"engine": "python-docx", "paras": len(document.paragraphs)},
                }
            except Exception as exc:
                return {"ok": False, "text": "", "meta": {"error": f"DOCX parse failed: {exc}"}}

        system = platform.system().lower()
        if suffix == ".doc" and system.startswith("win") and win32com:
            com_initialized = False
            try:
                if pythoncom:
                    pythoncom.CoInitialize()
                    com_initialized = True
                word = win32com.Dispatch("Word.Application")
                word.Visible = False
                try:
                    doc_obj = word.Documents.Open(str(p), ReadOnly=True)
                    try:
                        text = doc_obj.Content.Text or ""
                    finally:
                        doc_obj.Close(False)
                finally:
                    try:
                        word.Quit()
                    except Exception:
                        pass
                return {
                    "ok": True,
                    "text": TextCleaner.clean(text),
                    "meta": {"engine": "win32com-word"},
                }
            except Exception as exc:
                return {"ok": False, "text": "", "meta": {"error": f"DOC win32com 실패: {exc}"}}
            finally:
                if com_initialized and pythoncom:
                    try:
                        pythoncom.CoUninitialize()
                    except Exception:
                        pass

        return {
            "ok": False,
            "text": "",
            "meta": {"error": "DOC/DOCX 추출을 위해 python-docx 또는 Windows Word가 필요합니다."},
        }

class ExcelLikeExtractor(BaseExtractor):
    exts=(".xlsx",".xls",".xlsm",".xlsb",".xltx",".csv")
    def extract(self, p:Path)->Dict[str,Any]:
        if pd is None:
            return {"ok":False,"text":"","meta":{"error":"pandas required"}}
        try:
            max_bytes_env = os.getenv("INFOPILOT_EXCEL_MAX_BYTES", "").strip()
            max_bytes = int(max_bytes_env) if max_bytes_env else 25 * 1024 * 1024
            try:
                size = p.stat().st_size
            except Exception:
                size = None
            if size is not None and size > max_bytes:
                return {
                    "ok": False,
                    "text": "",
                    "meta": {"error": "excel/csv file too large", "size": size, "max_bytes": max_bytes},
                }

            if p.suffix.lower()==".csv":
                df=pd.read_csv(p, nrows=200, encoding="utf-8", engine="python")
                txt=self._df_to_text(df)
                return {"ok":True,"text":txt,"meta":{"engine":"pandas","columns":df.columns.tolist(), "rows_preview":min(200,len(df))}}
            eng = "openpyxl" if p.suffix.lower() in (".xlsx",".xlsm",".xltx") else ("xlrd" if p.suffix.lower()==".xls" else "pyxlsb")

            max_sheets_env = os.getenv("INFOPILOT_EXCEL_MAX_SHEETS", "").strip()
            max_sheets = int(max_sheets_env) if max_sheets_env else 3
            nrows_env = os.getenv("INFOPILOT_EXCEL_NROWS", "").strip()
            nrows = int(nrows_env) if nrows_env else 200

            parts=[]
            sheet_names: List[str] = []
            with pd.ExcelFile(p, engine=eng) as xf:
                sheet_names = list(getattr(xf, "sheet_names", []) or [])
                if not sheet_names:
                    sheet_names = [0]  # type: ignore[list-item]
                for sheet in sheet_names[: max(1, max_sheets)]:
                    df_sheet = xf.parse(sheet, nrows=nrows)
                    sheet_label = sheet if isinstance(sheet, str) else str(sheet)
                    parts.append(f"[Sheet:{sheet_label}]")
                    parts.append(" | ".join(map(str, df_sheet.columns.tolist())))
                    for _,row in df_sheet.head(50).iterrows():
                        parts.append(" • "+" | ".join(map(lambda x: str(x), row.tolist())))

            meta = {"engine":"pandas","sheets":sheet_names}
            if len(sheet_names) > max(1, max_sheets):
                meta["sheets_truncated"] = True
                meta["sheets_kept"] = sheet_names[: max(1, max_sheets)]
            meta["rows_preview"] = nrows
            return {"ok":True,"text":TextCleaner.clean("\n".join(parts)),"meta":meta}
        except Exception as e:
            detail = str(e)
            if "openpyxl" in detail.lower():
                detail += " (pip install openpyxl)"
            return {"ok":False,"text":"","meta":{"error":f"excel/csv read failed: {detail}"}}
    @staticmethod
    def _df_to_text(df)->str:
        cols=" | ".join(map(str, df.columns.tolist()))
        rows=[]
        for _,row in df.head(50).iterrows():
            rows.append(" • "+" | ".join(map(lambda x: str(x), row.tolist())))
        return TextCleaner.clean(f"{cols}\n"+"\n".join(rows))

class PdfExtractor(BaseExtractor):
    exts = (".pdf",)

    def extract(self, p: Path) -> Dict[str, Any]:
        max_pages_env = os.getenv("INFOPILOT_PDF_MAX_PAGES", "").strip()
        max_pages = int(max_pages_env) if max_pages_env else 200
        max_chars_env = os.getenv("INFOPILOT_PDF_MAX_CHARS", "").strip()
        max_chars = int(max_chars_env) if max_chars_env else 200_000

        if fitz:
            try:
                with fitz.open(str(p)) as doc:
                    page_count = doc.page_count
                    parts: List[str] = []
                    truncated_pages = False
                    char_count = 0
                    for idx, page in enumerate(doc):
                        if idx >= max_pages:
                            truncated_pages = True
                            break
                        page_text = page.get_text("text") or ""
                        parts.append(page_text)
                        char_count += len(page_text)
                        if char_count >= max_chars:
                            truncated_pages = True
                            break
                    text = "\n".join(parts)
                return {
                    "ok": True,
                    "text": TextCleaner.clean(text),
                    "meta": {
                        "engine": "pymupdf",
                        "pages": page_count,
                        "max_pages": max_pages,
                        "max_chars": max_chars,
                        "truncated": bool(truncated_pages),
                    },
                }
            except Exception:
                pass
        if pdfplumber:
            try:
                with pdfplumber.open(str(p)) as doc:
                    pages: List[str] = []
                    truncated_pages = False
                    char_count = 0
                    for idx, page in enumerate(doc.pages):
                        if idx >= max_pages:
                            truncated_pages = True
                            break
                        page_text = page.extract_text() or ""
                        pages.append(page_text)
                        char_count += len(page_text)
                        if char_count >= max_chars:
                            truncated_pages = True
                            break
                text = "\n".join(pages)
                cleaned = TextCleaner.clean(text)
                if cleaned:
                    return {
                        "ok": True,
                        "text": cleaned,
                        "meta": {
                            "engine": "pdfplumber",
                            "pages": len(pages),
                            "max_pages": max_pages,
                            "max_chars": max_chars,
                            "truncated": bool(truncated_pages),
                        },
                    }
            except Exception as exc:
                return {"ok": False, "text": "", "meta": {"error": f"PDF pdfplumber 실패: {exc}"}}
        if pdfminer_extract_text:
            try:
                text = pdfminer_extract_text(str(p))
                return {"ok": True, "text": TextCleaner.clean(text), "meta": {"engine": "pdfminer"}}
            except Exception as exc:
                return {"ok": False, "text": "", "meta": {"error": f"PDF pdfminer 실패: {exc}"}}
        return {"ok": False, "text": "", "meta": {"error": "PDF 추출 엔진이 설치되지 않았습니다."}}


class PptExtractor(BaseExtractor):
    exts = (".ppt", ".pptx")

    def extract(self, p: Path) -> Dict[str, Any]:
        suffix = p.suffix.lower()
        if suffix == ".pptx" and pptx:
            try:
                presentation = pptx.Presentation(str(p))
                texts: List[str] = []
                for idx, slide in enumerate(presentation.slides, 1):
                    parts: List[str] = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text = (shape.text or "").strip()
                            if text:
                                parts.append(text)
                    if parts:
                        texts.append(f"[Slide {idx}] " + " ".join(parts))
                return {
                    "ok": True,
                    "text": TextCleaner.clean("\n".join(texts)),
                    "meta": {"engine": "python-pptx", "slides": len(presentation.slides)},
                }
            except Exception as exc:
                return {"ok": False, "text": "", "meta": {"error": f"PPTX parse failed: {exc}"}}

        system = platform.system().lower()
        if suffix == ".ppt" and system.startswith("win") and win32com:
            com_initialized = False
            try:
                if pythoncom:
                    pythoncom.CoInitialize()
                    com_initialized = True
                powerpoint = win32com.Dispatch("PowerPoint.Application")
                powerpoint.Visible = False
                presentation = powerpoint.Presentations.Open(str(p), WithWindow=False)
                texts: List[str] = []
                try:
                    for slide in presentation.Slides:
                        parts = []
                        for shape in slide.Shapes:
                            has_text = hasattr(shape, "HasTextFrame") and shape.HasTextFrame
                            if has_text and shape.TextFrame.HasText:
                                parts.append(shape.TextFrame.TextRange.Text)
                        if parts:
                            texts.append(" ".join(parts))
                    return {
                        "ok": True,
                        "text": TextCleaner.clean("\n".join(texts)),
                        "meta": {"engine": "win32com-ppt"},
                    }
                finally:
                    presentation.Close()
                    powerpoint.Quit()
            except Exception as exc:
                return {"ok": False, "text": "", "meta": {"error": f"PPT win32com 실패: {exc}"}}
            finally:
                if com_initialized and pythoncom:
                    try:
                        pythoncom.CoUninitialize()
                    except Exception:
                        pass

        return {"ok": False, "text": "", "meta": {"error": "PPT/PPTX 추출을 위해 python-pptx 또는 Windows PowerPoint가 필요합니다."}}


class PlainTextExtractor(BaseExtractor):
    exts = (".txt", ".md", ".rst", ".log")

    def extract(self, p: Path) -> Dict[str, Any]:
        try:
            raw_text = p.read_text(encoding="utf-8", errors="replace")
        except Exception as exc:
            return {"ok": False, "text": "", "meta": {"error": f"텍스트 파일 읽기 실패: {exc}"}}
        cleaned = TextCleaner.clean(raw_text)
        meta: Dict[str, Any] = {"engine": "plain-text"}
        if p.suffix.lower() == ".md":
            meta["format"] = "markdown"
        return {
            "ok": bool(cleaned),
            "text": cleaned,
            "text_original": raw_text,
            "meta": meta,
        }


class CodeExtractor(BaseExtractor):
    exts = (
        ".py",
        ".json",
        ".yaml",
        ".yml",
        ".toml",
        ".ini",
        ".cfg",
        ".sh",
        ".bash",
    )

    def extract(self, p: Path) -> Dict[str, Any]:
        try:
            raw_text = p.read_text(encoding="utf-8", errors="replace")
        except Exception as exc:
            return {"ok": False, "text": "", "meta": {"error": f"코드 파일 읽기 실패: {exc}"}}
        cleaned = TextCleaner.clean(raw_text)
        meta: Dict[str, Any] = {"engine": "code", "extension": p.suffix.lower()}
        return {
            "ok": bool(cleaned),
            "text": cleaned,
            "text_original": raw_text,
            "meta": meta,
        }


EXTRACTORS = [
    HwpExtractor(),
    DocDocxExtractor(),
    ExcelLikeExtractor(),
    PdfExtractor(),
    PptExtractor(),
    PlainTextExtractor(),
    CodeExtractor(),
]
EXT_MAP={e:ex for ex in EXTRACTORS for e in ex.exts}
